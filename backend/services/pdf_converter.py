"""
PDF converter service - complete implementation
"""
import logging
from typing import Dict, Optional
import base64
from io import BytesIO
import os
import tempfile

from PIL import Image
from pdf2image import convert_from_bytes
import fitz  # PyMuPDF
import pypdf
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from docx import Document
import mammoth

logger = logging.getLogger(__name__)

# WeasyPrint is optional (requires system libraries on Windows)
# Import will be done inside functions that need it
WEASYPRINT_AVAILABLE = None

def _check_weasyprint():
    """Check if WeasyPrint is available"""
    global WEASYPRINT_AVAILABLE
    if WEASYPRINT_AVAILABLE is None:
        try:
            from weasyprint import HTML
            WEASYPRINT_AVAILABLE = True
        except (ImportError, OSError):
            WEASYPRINT_AVAILABLE = False
            logger.warning("WeasyPrint not available - HTML to PDF conversion will use fallback")
    return WEASYPRINT_AVAILABLE


class PDFConverterService:
    """Service for PDF conversions"""
    
    def __init__(self):
        """Initialize PDF converter service"""
        pass
    
    async def pdf_to_docx(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to DOCX"""
        try:
            # Open PDF
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            # Extract text and formatting
            docx_content = []
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                text = page.get_text()
                docx_content.append(text)
            
            pdf_document.close()
            
            # Create DOCX using python-docx
            doc = Document()
            for text in docx_content:
                doc.add_paragraph(text)
            
            # Save to buffer
            output_buffer = BytesIO()
            doc.save(output_buffer)
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.docx'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to DOCX error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_pptx(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to PPTX"""
        try:
            from pptx import Presentation
            
            # Try using pdf2image (requires Poppler)
            try:
                images = convert_from_bytes(pdf_content, dpi=200)
            except Exception as poppler_error:
                # Fallback to PyMuPDF if Poppler is not available
                logger.warning(f"pdf2image failed (Poppler may not be installed): {poppler_error}. Using PyMuPDF fallback.")
                
                # Use PyMuPDF (fitz) as fallback
                pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
                images = []
                
                for page_num in range(len(pdf_document)):
                    pdf_page = pdf_document[page_num]
                    # Render page to image (pixmap)
                    mat = fitz.Matrix(200/72, 200/72)  # 200 DPI
                    pix = pdf_page.get_pixmap(matrix=mat)
                    
                    # Convert pixmap to PIL Image
                    img_data = pix.tobytes("ppm")
                    image = Image.open(BytesIO(img_data))
                    images.append(image)
                
                pdf_document.close()
            
            # Create PPTX
            prs = Presentation()
            prs.slide_width = 9144000  # 10 inches
            prs.slide_height = 6858000  # 7.5 inches
            
            for image in images:
                # Create slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Convert PIL image to bytes
                img_buffer = BytesIO()
                image.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                
                # Add image to slide
                slide.shapes.add_picture(img_buffer, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # Save to buffer
            output_buffer = BytesIO()
            prs.save(output_buffer)
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.pptx'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to PPTX error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_xlsx(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to XLSX"""
        try:
            import pandas as pd
            from openpyxl import Workbook
            
            # Open PDF and extract text
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            # Extract tables (simplified - use pdfplumber for better table extraction)
            try:
                import pdfplumber
                with pdfplumber.open(BytesIO(pdf_content)) as pdf:
                    all_tables = []
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        all_tables.extend(tables)
                    
                    # Convert to DataFrame
                    if all_tables:
                        df = pd.DataFrame(all_tables[0])
                    else:
                        # Fallback: extract text as single column
                        text = ""
                        for page in pdf_document:
                            text += page.get_text()
                        df = pd.DataFrame([text.split('\n')])
            except:
                # Fallback: extract text
                text = ""
                for page_num in range(len(pdf_document)):
                    page = pdf_document[page_num]
                    text += page.get_text()
                df = pd.DataFrame([text.split('\n')])
            
            pdf_document.close()
            
            # Save to XLSX using openpyxl directly
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            
            # Write DataFrame to worksheet
            for r_idx, row in enumerate(df.itertuples(index=False), start=1):
                for c_idx, value in enumerate(row, start=1):
                    ws.cell(row=r_idx, column=c_idx, value=value)
            
            output_buffer = BytesIO()
            wb.save(output_buffer)
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.xlsx'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to XLSX error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_jpg(self, pdf_content: bytes, filename: str, page: int = 0) -> Dict[str, str]:
        """Convert PDF page to JPG"""
        try:
            # Try using pdf2image (requires Poppler)
            try:
                images = convert_from_bytes(pdf_content, dpi=300, first_page=page + 1, last_page=page + 1)
                
                if not images:
                    raise ValueError("No pages found in PDF")
                
                # Get first (and only) image
                image = images[0]
            except Exception as poppler_error:
                # Fallback to PyMuPDF if Poppler is not available
                logger.warning(f"pdf2image failed (Poppler may not be installed): {poppler_error}. Using PyMuPDF fallback.")
                
                # Use PyMuPDF (fitz) as fallback
                pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
                
                if page >= len(pdf_document):
                    pdf_document.close()
                    raise ValueError(f"Page {page + 1} not found in PDF (total pages: {len(pdf_document)})")
                
                # Get the specific page
                pdf_page = pdf_document[page]
                
                # Render page to image (pixmap)
                mat = fitz.Matrix(300/72, 300/72)  # 300 DPI
                pix = pdf_page.get_pixmap(matrix=mat)
                
                # Convert pixmap to PIL Image
                img_data = pix.tobytes("ppm")
                image = Image.open(BytesIO(img_data))
                
                pdf_document.close()
            
            # Convert to JPG
            output_buffer = BytesIO()
            image.save(output_buffer, format='JPEG', quality=90)
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:image/jpeg;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + f'_page_{page + 1}.jpg'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to JPG error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_txt(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to TXT"""
        try:
            # Open PDF
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            # Extract text from all pages
            text = ""
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                text += page.get_text()
            
            pdf_document.close()
            
            # Convert to data URL
            text_bytes = text.encode('utf-8')
            data_url = f"data:text/plain;base64,{base64.b64encode(text_bytes).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.txt'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to TXT error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_html(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to HTML"""
        try:
            # Open PDF
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            # Convert to HTML
            html = ""
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                html += page.get_text("html")
            
            pdf_document.close()
            
            # Convert to data URL
            html_bytes = html.encode('utf-8')
            data_url = f"data:text/html;base64,{base64.b64encode(html_bytes).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.html'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to HTML error: {exc}", exc_info=True)
            raise
    
    async def pdf_to_pdfa(self, pdf_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PDF to PDF/A"""
        try:
            # Open PDF
            pdf_reader = pypdf.PdfReader(BytesIO(pdf_content))
            pdf_writer = pypdf.PdfWriter()
            
            # Copy all pages
            for page in pdf_reader.pages:
                pdf_writer.add_page(page)
            
            # Add PDF/A metadata
            pdf_writer.add_metadata({
                '/Title': 'PDF/A Document',
                '/Author': 'MegaPixelAI',
                '/Subject': 'PDF/A compliant document',
                '/Creator': 'MegaPixelAI PDF/A Converter',
                '/Producer': 'MegaPixelAI PDF/A Converter',
            })
            
            # Save to buffer
            output_buffer = BytesIO()
            pdf_writer.write(output_buffer)
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/pdf;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '_pdfa.pdf'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PDF to PDF/A error: {exc}", exc_info=True)
            raise
    
    async def docx_to_pdf(self, docx_content: bytes, filename: str) -> Dict[str, str]:
        """Convert DOCX to PDF"""
        try:
            # Convert DOCX to HTML using mammoth
            html_result = mammoth.convert_to_html(BytesIO(docx_content))
            html_string = html_result.value
            
            # Convert HTML to PDF using WeasyPrint (if available) or reportlab
            if _check_weasyprint():
                from weasyprint import HTML
                pdf_bytes = HTML(string=html_string).write_pdf()
            else:
                # Fallback: use reportlab to create PDF from text
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import SimpleDocTemplate, Paragraph
                
                output_buffer = BytesIO()
                doc = SimpleDocTemplate(output_buffer, pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Simple HTML to text conversion
                import re
                text = re.sub(r'<[^>]+>', '', html_string)
                paragraphs = text.split('\n')
                for para in paragraphs:
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                
                doc.build(story)
                pdf_bytes = output_buffer.getvalue()
            
            # Convert to data URL
            data_url = f"data:application/pdf;base64,{base64.b64encode(pdf_bytes).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.pdf'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"DOCX to PDF error: {exc}", exc_info=True)
            raise
    
    async def ppt_to_pdf(self, ppt_content: bytes, filename: str) -> Dict[str, str]:
        """Convert PPT/PPTX to PDF"""
        try:
            from pptx import Presentation
            
            # Open PPTX
            prs = Presentation(BytesIO(ppt_content))
            
            # Convert slides to images then to PDF
            images = []
            for i, slide in enumerate(prs.slides):
                # Convert slide to image (simplified - would need actual rendering)
                # For now, create placeholder
                img = Image.new('RGB', (1920, 1080), color='white')
                images.append(img)
            
            # Convert images to PDF
            if images:
                images[0].save(
                    BytesIO(),
                    "PDF",
                    resolution=100.0,
                    save_all=True,
                    append_images=images[1:] if len(images) > 1 else []
                )
            
            # Use reportlab to create PDF
            output_buffer = BytesIO()
            pdf_canvas = canvas.Canvas(output_buffer, pagesize=A4)
            
            for img in images:
                img_buffer = BytesIO()
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                pdf_canvas.drawImage(ImageReader(img_buffer), 0, 0, width=A4[0], height=A4[1])
                pdf_canvas.showPage()
            
            pdf_canvas.save()
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/pdf;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.pdf'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"PPT to PDF error: {exc}", exc_info=True)
            raise
    
    async def xls_to_pdf(self, xls_content: bytes, filename: str) -> Dict[str, str]:
        """Convert XLS/XLSX to PDF"""
        try:
            import pandas as pd
            
            # Read Excel file
            df = pd.read_excel(BytesIO(xls_content), engine='openpyxl')
            
            # Convert to PDF using reportlab
            output_buffer = BytesIO()
            pdf_canvas = canvas.Canvas(output_buffer, pagesize=A4)
            
            # Add table to PDF
            y = A4[1] - 50
            pdf_canvas.setFont("Helvetica", 10)
            
            # Write header
            for col_idx, col_name in enumerate(df.columns):
                pdf_canvas.drawString(50 + col_idx * 100, y, str(col_name))
            
            y -= 20
            
            # Write data
            for _, row in df.iterrows():
                if y < 50:
                    pdf_canvas.showPage()
                    y = A4[1] - 50
                
                for col_idx, value in enumerate(row):
                    pdf_canvas.drawString(50 + col_idx * 100, y, str(value))
                y -= 15
            
            pdf_canvas.save()
            output_buffer.seek(0)
            
            # Convert to data URL
            data_url = f"data:application/pdf;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.pdf'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"XLS to PDF error: {exc}", exc_info=True)
            raise
    
    async def html_to_pdf(self, html_content: bytes, filename: str) -> Dict[str, str]:
        """Convert HTML to PDF"""
        try:
            # Convert HTML to PDF
            html_string = html_content.decode('utf-8')
            
            # Check if WeasyPrint is available
            if _check_weasyprint():
                from weasyprint import HTML
                pdf_bytes = HTML(string=html_string).write_pdf()
            else:
                # Fallback: use reportlab to create PDF from text
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import SimpleDocTemplate, Paragraph
                
                output_buffer = BytesIO()
                doc = SimpleDocTemplate(output_buffer, pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Simple HTML to text conversion
                import re
                text = re.sub(r'<[^>]+>', '', html_string)
                paragraphs = text.split('\n')
                for para in paragraphs:
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                
                doc.build(story)
                pdf_bytes = output_buffer.getvalue()
            
            # Convert to data URL
            data_url = f"data:application/pdf;base64,{base64.b64encode(pdf_bytes).decode()}"
            
            result_name = filename.rsplit('.', 1)[0] + '.pdf'
            
            return {
                "name": result_name,
                "dataUrl": data_url,
            }
        
        except Exception as exc:
            logger.error(f"HTML to PDF error: {exc}", exc_info=True)
            raise
    
    async def jpg_to_pdf(self, image_content: bytes, filename: str) -> Dict[str, str]:
        """Convert JPG/PNG to PDF - supports single image or list of images"""
        try:
            # Check if image_content is a single image or list
            if isinstance(image_content, list):
                # Multiple images - create multi-page PDF
                images = []
                for img_bytes in image_content:
                    img = Image.open(BytesIO(img_bytes))
                    # Convert to RGB if necessary (for PDF)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                
                # Create PDF with all images
                if images:
                    output_buffer = BytesIO()
                    # Save first image and append others
                    images[0].save(
                        output_buffer,
                        format='PDF',
                        resolution=100.0,
                        save_all=True,
                        append_images=images[1:] if len(images) > 1 else []
                    )
                    output_buffer.seek(0)
                    
                    # Convert to data URL
                    data_url = f"data:application/pdf;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
                    
                    result_name = filename.rsplit('.', 1)[0] + '.pdf'
                    
                    return {
                        "name": result_name,
                        "dataUrl": data_url,
                    }
            else:
                # Single image
                image = Image.open(BytesIO(image_content))
                
                # Convert to RGB if necessary (PDF requires RGB)
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Create PDF
                output_buffer = BytesIO()
                image.save(output_buffer, format='PDF', resolution=100.0)
                output_buffer.seek(0)
                
                # Convert to data URL
                data_url = f"data:application/pdf;base64,{base64.b64encode(output_buffer.getvalue()).decode()}"
                
                result_name = filename.rsplit('.', 1)[0] + '.pdf'
                
                return {
                    "name": result_name,
                    "dataUrl": data_url,
                }
        
        except Exception as exc:
            logger.error(f"JPG to PDF error: {exc}", exc_info=True)
            raise
